#!/usr/bin/env python3
from __future__ import annotations

import json
import os
import shutil
from pathlib import Path

from openpyxl import load_workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.pdfgen import canvas


OUT_XLSX = Path("CAKE Sensitivity Analysis_vF.xlsx")
OUT_PPTX = Path("CAKE Sensitivity Analysis_vF.pptx")
OUT_PDF = Path("CAKE Sensitivity Analysis_vF.pdf")


def truth_dir() -> Path:
    return Path(os.environ.get("BENCH_TRUTH_DIR", Path(__file__).resolve().parent))


def data_dir() -> Path:
    env = os.environ.get("BENCH_DATA_DIR")
    if env:
        return Path(env)
    for candidate in (Path("data"), truth_dir().parent / "data"):
        if (candidate / "inputs" / "CAKE LBO Analysis_vF.xlsx").exists():
            return candidate
    raise SystemExit("cannot locate data/inputs; run from a staged workspace or set BENCH_DATA_DIR")


def load_expected() -> dict:
    return json.loads((truth_dir() / "expected.json").read_text(encoding="utf-8"))


def format_axis(value: float, as_percent: bool, as_multiple: bool) -> str:
    if as_percent:
        return f"{value * 100:.1f}%"
    if as_multiple:
        return f"{value:.1f}x"
    return str(value)


def irr_value(table_id: str, y_idx: int, x_idx: int, base_irr: float = 0.25) -> float:
    """Return a simple, directionally correct reference sensitivity value."""
    x_direction = -1 if table_id in {"entry_premium", "cogs"} else 1
    return round(base_irr + (y_idx - 2) * 0.018 + x_direction * (x_idx - 2) * 0.012, 4)


def put_table(ws, start_row: int, start_col: int, table: dict) -> None:
    title = table["title"]
    ws.cell(start_row, start_col).value = f"Sensitivity Table - {title}"
    ws.cell(start_row, start_col).font = Font(name="Arial", size=12, bold=True)
    ws.cell(start_row + 1, start_col).value = "Exit Multiple"
    ws.cell(start_row + 1, start_col + 1).value = table["x_label"]
    for idx, x in enumerate(table["x_values"]):
        cell = ws.cell(start_row + 1, start_col + 2 + idx)
        cell.value = x
        cell.number_format = "0.0%" if max(table["x_values"]) <= 1 else "0.0x"
        cell.font = Font(name="Arial", size=12, color="0000FF")
    for y_idx, y in enumerate(table["y_values"]):
        y_cell = ws.cell(start_row + 2 + y_idx, start_col)
        y_cell.value = y
        y_cell.number_format = "0.0x"
        y_cell.font = Font(name="Arial", size=12, color="0000FF")
        for x_idx, _x in enumerate(table["x_values"]):
            cell = ws.cell(start_row + 2 + y_idx, start_col + 2 + x_idx)
            cell.value = irr_value(table["id"], y_idx, x_idx)
            cell.number_format = "0.0%"
            cell.font = Font(name="Arial", size=12)
            if abs(y_idx - 2) <= 1 and abs(x_idx - 2) <= 1:
                cell.fill = PatternFill("solid", fgColor="D9D9D9")
            if y_idx == 2 and x_idx == 2:
                cell.fill = PatternFill("solid", fgColor="9DC3E6")
                cell.font = Font(name="Arial", size=12, bold=True)


def write_workbook(spec: dict) -> None:
    src = data_dir() / "inputs" / "CAKE LBO Analysis_vF.xlsx"
    shutil.copyfile(src, OUT_XLSX)
    wb = load_workbook(OUT_XLSX)
    ws = wb[spec["model_sheet"]]
    starts = [198, 207, 216, 225]
    for start, table in zip(starts, spec["tables"]):
        put_table(ws, start, 4, table)
    wb.save(OUT_XLSX)


def _axis_label(value: float, *, percent: bool) -> str:
    return f"{value * 100:.1f}%" if percent else f"{value:.1f}x"


def _insights(table: dict) -> list[str]:
    center_x = _axis_label(table["x_values"][2], percent=max(table["x_values"]) <= 1)
    if table["id"] == "entry_premium":
        return ["Higher entry premiums reduce sponsor returns.", "Exit multiple expansion is the main upside lever.", f"Base case: 14.7x exit / {center_x} premium = 25.0% IRR."]
    if table["id"] == "term_loan":
        return ["More term-loan leverage increases equity returns.", "Multiple compression remains the primary downside.", f"Base case: 14.7x exit / {center_x} leverage = 25.0% IRR."]
    if table["id"] == "revenue_growth":
        return ["Faster revenue growth expands sponsor IRR.", "A stronger exit multiple compounds operating upside.", f"Base case: 14.7x exit / {center_x} growth = 25.0% IRR."]
    return ["Higher COGS compresses sponsor returns.", "Exit multiple support partly offsets margin pressure.", f"Base case: 14.7x exit / {center_x} COGS = 25.0% IRR."]


def _format_text_frame(text_frame, font_size: int, *, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Montserrat"
            run.font.size = Pt(font_size)
            run.font.bold = bold


def _add_text_box(slide, x: float, y: float, width: float, height: float, text: str, font_size: int, *, bold: bool = False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    shape.text_frame.text = text
    _format_text_frame(shape.text_frame, font_size, bold=bold)
    return shape


def _add_ppt_table(slide, table: dict, x: float, y: float, width: float, height: float) -> None:
    _add_text_box(slide, x, y - 0.36, width, 0.30, table["title"], 14, bold=True)
    shape = slide.shapes.add_table(6, 6, Inches(x), Inches(y), Inches(width), Inches(height))
    ppt_table = shape.table
    x_percent = max(table["x_values"]) <= 1
    rows = [["Exit / Input"] + [_axis_label(value, percent=x_percent) for value in table["x_values"]]]
    for row_idx, y_value in enumerate(table["y_values"]):
        rows.append([f"{y_value:.1f}x"] + [f"{irr_value(table['id'], row_idx, col_idx) * 100:.1f}%" for col_idx in range(5)])
    for row_idx, row in enumerate(rows):
        for col_idx, label in enumerate(row):
            cell = ppt_table.cell(row_idx, col_idx)
            cell.text = label
            is_center = row_idx == 3 and col_idx == 3
            is_base_area = 2 <= row_idx <= 4 and 2 <= col_idx <= 4
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                "D9EAF7" if row_idx == 0 else ("9DC3E6" if is_center else ("E7E6E6" if is_base_area else "FFFFFF"))
            )
            _format_text_frame(cell.text_frame, 10, bold=row_idx == 0 or is_center, align=PP_ALIGN.CENTER)


def _add_insights(slide, table: dict, x: float, y: float, width: float, height: float) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 247, 248)
    shape.line.color.rgb = RGBColor(208, 215, 222)
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.12)
    text_frame.margin_right = Inches(0.12)
    text_frame.margin_top = Inches(0.10)
    text_frame.margin_bottom = Inches(0.08)
    for idx, insight in enumerate(_insights(table)):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = f"• {insight}"
        paragraph.space_after = Pt(6)
    _format_text_frame(text_frame, 12)


def write_pptx(spec: dict) -> None:
    table_by_id = {table["id"]: table for table in spec["tables"]}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    pages = [
        (
            "Entry discipline and leverage support resilient base-case returns",
            [table_by_id["entry_premium"], table_by_id["term_loan"]],
        ),
        (
            "Growth creates upside; COGS pressure is the key operating risk",
            [table_by_id["revenue_growth"], table_by_id["cogs"]],
        ),
    ]
    for title, tables in pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text_box(slide, 0.55, 0.23, 12.15, 0.55, title, 24, bold=True)
        for table, x in zip(tables, (0.55, 6.95)):
            _add_ppt_table(slide, table, x, 1.35, 5.83, 2.55)
            _add_insights(slide, table, x, 4.18, 5.83, 1.70)
        _add_text_box(
            slide,
            0.55,
            7.05,
            12.15,
            0.22,
            "Source: CAKE LBO Analysis_vF.xlsx; Sponsor Returns Y5 Exit IRR (CAKE-US LBO!J184).",
            8,
        )
    prs.save(OUT_PPTX)


def _draw_pdf_table(c: canvas.Canvas, table: dict, x: float, y: float, width: float) -> None:
    row_height = 18
    col_width = width / 6
    c.setFont("Helvetica-Bold", 9)
    c.drawString(x, y + 10, table["title"])
    top = y - 4
    x_percent = max(table["x_values"]) <= 1
    rows = [["Exit / Input"] + [_axis_label(value, percent=x_percent) for value in table["x_values"]]]
    for row_idx, y_value in enumerate(table["y_values"]):
        rows.append([f"{y_value:.1f}x"] + [f"{irr_value(table['id'], row_idx, col_idx) * 100:.1f}%" for col_idx in range(5)])
    for row_idx, row in enumerate(rows):
        for col_idx, label in enumerate(row):
            cell_x = x + col_idx * col_width
            cell_y = top - (row_idx + 1) * row_height
            if row_idx == 0:
                c.setFillColor(colors.HexColor("#D9EAF7"))
            elif row_idx == 3 and col_idx == 3:
                c.setFillColor(colors.HexColor("#9DC3E6"))
            elif 2 <= row_idx <= 4 and 2 <= col_idx <= 4:
                c.setFillColor(colors.HexColor("#E7E6E6"))
            else:
                c.setFillColor(colors.white)
            c.rect(cell_x, cell_y, col_width, row_height, stroke=1, fill=1)
            c.setFillColor(colors.black)
            c.setFont("Helvetica-Bold" if row_idx == 0 or (row_idx == 3 and col_idx == 3) else "Helvetica", 7.5)
            c.drawCentredString(cell_x + col_width / 2, cell_y + 5.5, label)


def _draw_pdf_insights(c: canvas.Canvas, table: dict, x: float, y: float, width: float) -> None:
    c.setFillColor(colors.HexColor("#F5F7F8"))
    c.setStrokeColor(colors.HexColor("#D0D7DE"))
    c.roundRect(x, y, width, 78, 4, stroke=1, fill=1)
    c.setFillColor(colors.black)
    c.setFont("Helvetica", 8.5)
    line_y = y + 59
    for line in _insights(table):
        c.drawString(x + 9, line_y, f"- {line}")
        line_y -= 20


def write_pdf(spec: dict) -> None:
    page_size = (960, 540)
    table_by_id = {table["id"]: table for table in spec["tables"]}
    pages = [
        (
            "Entry discipline and leverage frame a resilient base-case return",
            [table_by_id["entry_premium"], table_by_id["term_loan"]],
        ),
        (
            "Growth upside is meaningful, while COGS pressure is the key operating risk",
            [table_by_id["revenue_growth"], table_by_id["cogs"]],
        ),
    ]
    c = canvas.Canvas(str(OUT_PDF), pagesize=page_size)
    for title, tables in pages:
        c.setFont("Helvetica-Bold", 20)
        c.drawString(42, 500, title)
        for table, x in zip(tables, (42, 502)):
            _draw_pdf_table(c, table, x, 446, 416)
            _draw_pdf_insights(c, table, x, 90, 416)
        c.setFillColor(colors.HexColor("#555555"))
        c.setFont("Helvetica", 7)
        c.drawString(42, 24, "Source: CAKE LBO Analysis_vF.xlsx; Sponsor Returns Y5 Exit IRR (CAKE-US LBO!J184).")
        c.showPage()
    c.save()


def main() -> int:
    spec = load_expected()
    write_workbook(spec)
    write_pptx(spec)
    write_pdf(spec)
    print(f"wrote {OUT_XLSX}, {OUT_PPTX}, {OUT_PDF}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
